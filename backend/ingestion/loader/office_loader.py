from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

from backend.ingestion.loader.base import Loader


class WordLoader(Loader):
    async def load(self, file_path: str | Path) -> str:
        doc = DocxDocument(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def supported_extensions(self) -> set[str]:
        return {".docx"}


class ExcelLoader(Loader):
    async def load(self, file_path: str | Path) -> str:
        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                line = " | ".join(str(c) for c in row if c is not None)
                if line.strip():
                    lines.append(line)
        wb.close()
        return "\n".join(lines)

    def supported_extensions(self) -> set[str]:
        return {".xlsx", ".xls"}


class PPTLoader(Loader):
    async def load(self, file_path: str | Path) -> str:
        prs = Presentation(str(file_path))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)

    def supported_extensions(self) -> set[str]:
        return {".pptx"}
